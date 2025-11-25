from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import List, Dict


class PptxService:
    """Service for generating PowerPoint presentations"""
    
    def generate_presentation(self, title: str, slides: List[Dict[str, str]]) -> BytesIO:
        """
        Generate a .pptx file from slides
        
        Args:
            title: Presentation title
            slides: List of dicts with 'title' and 'content' keys
            
        Returns:
            BytesIO object containing the presentation
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add content slides
        for slide_data in slides:
            # Use bullet layout
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            
            # Add content
            if slide_data.get('content'):
                body_shape = slide.shapes.placeholders[1]
                text_frame = body_shape.text_frame
                text_frame.clear()
                
                # Parse bullet points
                lines = slide_data['content'].split('\n')
                for i, line in enumerate(lines):
                    line = line.strip()
                    if line:
                        # Remove existing bullet markers
                        line = line.lstrip('•-*').strip()
                        
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = line
                        p.level = 0
        
        # Save to BytesIO
        file_stream = BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        return file_stream


# Singleton instance
pptx_service = PptxService()
